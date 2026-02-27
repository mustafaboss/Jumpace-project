from pptx import Presentation
from pptx.util import Inches, Pt


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Construction Rules RAG Assistant"
    subtitle = slide.placeholders[1]
    subtitle.text = "PDF‑Grounded Chatbot for Field‑Ready Decisions"


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, line in enumerate(bullets):
        if i == 0:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = line
        p.level = 0


def build_presentation() -> None:
    prs = Presentation()

    # Slide 1 – Title
    add_title_slide(prs)

    # Slide 2 – Problem & Motivation
    add_bullet_slide(
        prs,
        "Problem & Motivation",
        [
            "Field teams struggle to quickly find construction rules for real scenarios (setbacks, signage, safety, etc.).",
            "Rules are buried in long PDFs and workers depend on verbal guidance or screenshots.",
            "Non‑compliance leads to delays, redesign, and safety risk.",
            "Goal: a fast, reliable assistant that answers strictly from official rule documents.",
        ],
    )

    # Slide 3 – Solution Overview
    add_bullet_slide(
        prs,
        "Solution Overview",
        [
            "A Retrieval Augmented Generation (RAG) chatbot focused on construction rules and regulations.",
            "Users describe a scenario in natural English; the assistant returns only the rules that apply.",
            "All answers are grounded in uploaded PDFs such as MUTCD or local authority rulebooks.",
            "If no rule is found in the PDF, the assistant clearly states that instead of guessing.",
        ],
    )

    # Slide 4 – Architecture (High Level)
    add_bullet_slide(
        prs,
        "High‑Level Architecture",
        [
            "Streamlit web UI for supervisors and engineers to chat with the assistant.",
            "Python backend using LangChain to orchestrate retrieval and generation.",
            "Pinecone vector database storing embeddings of all rulebook chunks with page metadata.",
            "OpenAI‑compatible models for embeddings and chat, with low temperature for precise answers.",
        ],
    )

    # Slide 5 – Implementation Journey
    add_bullet_slide(
        prs,
        "Implementation Journey",
        [
            "Set up a clean project structure with app, ingestion script, data folder, and requirements.",
            "Configured secure environment variables for all API keys and model names.",
            "Created a repeatable ingestion pipeline to load PDFs, chunk text, and index vectors in Pinecone.",
            "Built the RAG chain in LangChain and tuned prompts for safety and structure.",
            "Designed a modern, dashboard‑style Streamlit UI tailored for real construction workflows.",
        ],
    )

    # Slide 6 – Data Ingestion & Indexing
    add_bullet_slide(
        prs,
        "Data Ingestion & Indexing",
        [
            "Loaded rulebooks page‑by‑page using PyPDFLoader and attached file + page metadata.",
            "Chunked text with RecursiveCharacterTextSplitter (≈800–1200 tokens, overlapping context).",
            "Generated embeddings with text‑embedding‑3‑small for semantic search.",
            "Stored vectors in a Pinecone index (rag‑index) with cosine similarity for fast retrieval.",
        ],
    )

    # Slide 7 – Retrieval & Safety Controls
    add_bullet_slide(
        prs,
        "Retrieval & Safety Controls",
        [
            "Retriever always pulls context from the Pinecone index built over the PDFs.",
            "System prompt forbids external knowledge and guessing; answers must stay within context.",
            "If context does not contain an applicable rule, response is: \"This scenario’s relevant rule is not present in this PDF.\"",
            "Responses follow a fixed structure: scenario understanding, applicable rules, worker instructions, and page references.",
        ],
    )

    # Slide 8 – User Experience & UI
    add_bullet_slide(
        prs,
        "User Experience & UI",
        [
            "Streamlit interface styled as a modern OS‑like dashboard with gradient header and tabs.",
            "Dedicated chat panel with a single prominent call‑to‑action: Get Applicable Rules.",
            "Side panel with on‑site tips explaining how to phrase scenarios for best retrieval.",
            "Theme‑independent 3D/glass styling that looks premium on both light and dark modes.",
        ],
    )

    # Slide 9 – Demo Flow",
    add_bullet_slide(
        prs,
        "Demo Flow for the Client",
        [
            "Start with a realistic scenario (e.g., residential building on a main road‑facing plot).",
            "Show how the assistant explains the scenario, lists only relevant rules, and adds worker steps.",
            "Highlight page references so every answer is traceable back to the official rulebook.",
            "Demonstrate a case where no rule exists in the PDF to prove the system does not hallucinate.",
        ],
    )

    # Slide 10 – Benefits for the Client
    add_bullet_slide(
        prs,
        "Benefits for the Client",
        [
            "Faster, more confident decisions on site with immediate access to applicable rules.",
            "Reduced compliance risk because every answer can be traced to a specific PDF page.",
            "On‑site friendly instructions that translate regulations into concrete actions for workers.",
            "Scalable foundation that can support additional rulebooks, cities, and authorities over time.",
        ],
    )

    # Slide 11 – Future Enhancements
    add_bullet_slide(
        prs,
        "Future Enhancements",
        [
            "Support multiple authorities or jurisdictions with rulebook selection per project.",
            "Add user roles and an audit trail for compliance and internal QA.",
            "Allow engineers to give feedback on answers and continuously improve retrieval quality.",
            "Integrate with existing project management or permit systems for end‑to‑end workflows.",
        ],
    )

    prs.save("Construction_Rules_RAG_Assistant.pptx")


if __name__ == "__main__":
    build_presentation()

